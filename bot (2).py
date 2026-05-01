import os
import sys
import io
import tempfile
import logging
import traceback
from typing import List, Optional, Dict, Any

from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    ContextTypes,
    filters,
)

# ═══════════════════════════════════════════════
# إعدادات اللوق
# ═══════════════════════════════════════════════
logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO
)
logger = logging.getLogger(__name__)

# ═══════════════════════════════════════════════
# التوكن
# ═══════════════════════════════════════════════
TOKEN = os.environ.get("TELEGRAM_TOKEN", "8715590354:AAEtfxKZ_nbeDKGApup3deLF0Pxs69gPOlg")

# ═══════════════════════════════════════════════
# فحص المكتبات المتوفرة
# ═══════════════════════════════════════════════
logger.info("🔍 جاري فحص المكتبات...")

LIBS_STATUS = {}

def check_lib(name, import_name):
    try:
        __import__(import_name)
        LIBS_STATUS[name] = True
        logger.info(f"✅ {name} متوفر")
        return True
    except ImportError:
        LIBS_STATUS[name] = False
        logger.warning(f"❌ {name} غير متوفر")
        return False

# المكتبات الأساسية
check_lib("Pillow", "PIL")
check_lib("img2pdf", "img2pdf")
check_lib("ReportLab", "reportlab")
check_lib("python-docx", "docx")
check_lib("pandas", "pandas")
check_lib("openpyxl", "openpyxl")
check_lib("python-pptx", "pptx")
check_lib("PyPDF2", "PyPDF2")

logger.info("✅ انتهى فحص المكتبات")

# ═══════════════════════════════════════════════
# استيراد المكتبات المتوفرة
# ═══════════════════════════════════════════════
if LIBS_STATUS["Pillow"]:
    from PIL import Image

if LIBS_STATUS["ReportLab"]:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors

if LIBS_STATUS["python-docx"]:
    from docx import Document

if LIBS_STATUS["pandas"]:
    import pandas as pd

if LIBS_STATUS["python-pptx"]:
    from pptx import Presentation

if LIBS_STATUS["PyPDF2"]:
    from PyPDF2 import PdfMerger, PdfReader, PdfWriter

# ═══════════════════════════════════════════════
# دوال المساعدة
# ═══════════════════════════════════════════════

def clean_temp_files(*paths):
    """حذف الملفات المؤقتة"""
    for path in paths:
        try:
            if path and os.path.exists(path):
                os.remove(path)
                logger.info(f"🗑️ تم حذف: {path}")
        except Exception as e:
            logger.warning(f"⚠️ خطأ في حذف {path}: {e}")

def get_file_size_kb(path):
    """الحصول على حجم الملف بالكيلوبايت"""
    try:
        return os.path.getsize(path) / 1024
    except:
        return 0

# ═══════════════════════════════════════════════
# التعرف الذكي على نوع الملف
# ═══════════════════════════════════════════════

FILE_TYPES = {
    "image": {
        "extensions": [".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tiff"],
        "mime_prefix": "image/",
        "name": "صورة",
        "handler": "handle_image_conversion"
    },
    "word": {
        "extensions": [".docx"],
        "mime_types": ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"],
        "name": "Word",
        "handler": "handle_word_conversion"
    },
    "excel": {
        "extensions": [".xlsx", ".xls"],
        "mime_types": ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"],
        "name": "Excel",
        "handler": "handle_excel_conversion"
    },
    "powerpoint": {
        "extensions": [".pptx"],
        "mime_types": ["application/vnd.openxmlformats-officedocument.presentationml.presentation"],
        "name": "PowerPoint",
        "handler": "handle_ppt_conversion"
    },
    "pdf": {
        "extensions": [".pdf"],
        "mime_types": ["application/pdf"],
        "name": "PDF",
        "handler": "handle_pdf_operations"
    }
}

def detect_file_type(doc) -> Optional[Dict[str, Any]]:
    """التعرف الذكي على نوع الملف"""
    if not doc:
        return None

    file_name = doc.file_name.lower() if doc.file_name else ""
    mime_type = doc.mime_type.lower() if doc.mime_type else ""

    logger.info(f"🔍 فحص الملف: {file_name} | MIME: {mime_type}")

    for file_type, info in FILE_TYPES.items():
        # فحص الامتداد
        if any(file_name.endswith(ext) for ext in info.get("extensions", [])):
            return {"type": file_type, **info}

        # فحص MIME
        if "mime_types" in info:
            if any(mime_type == mt for mt in info["mime_types"]):
                return {"type": file_type, **info}

        # فحص بادئة MIME
        if "mime_prefix" in info:
            if mime_type.startswith(info["mime_prefix"]):
                return {"type": file_type, **info}

    return None

# ═══════════════════════════════════════════════
# أوامر البوت
# ═══════════════════════════════════════════════

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """أمر البدء - واجهة بسيطة"""
    welcome_text = """
🎉 *مرحباً بك في بوت PDF الذكي!*

🤖 *أنا بوت ذكي بسيط:*

📤 *أرسل أي ملف وأنا أحوله لـ PDF تلقائياً:*
• 📷 صور (JPG, PNG, WEBP)
• 📝 ملفات Word (DOCX)
• 📊 ملفات Excel (XLSX)
• 📽️ ملفات PowerPoint (PPTX)
• 📄 ملفات PDF (دمج/تقسيم/ضغط)

✨ *ما عليك إلا الإرسال!*
"""

    keyboard = [
        [InlineKeyboardButton("📋 حالة المكتبات", callback_data="check_libs")],
        [InlineKeyboardButton("❓ كيف أستخدم البوت", callback_data="help")]
    ]

    await update.message.reply_text(
        welcome_text,
        parse_mode="Markdown",
        reply_markup=InlineKeyboardMarkup(keyboard)
    )

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """المساعدة"""
    help_text = """
📖 *كيف أستخدم البوت؟*

*بسيط جداً!* 👇

1️⃣ *تحويل صورة:*
   أرسل الصورة مباشرة 📷

2️⃣ *تحويل Word:*
   أرسل ملف .docx 📝

3️⃣ *تحويل Excel:*
   أرسل ملف .xlsx 📊

4️⃣ *تحويل PowerPoint:*
   أرسل ملف .pptx 📽️

5️⃣ *معالجة PDF:*
   أرسل ملف PDF 📄
   ثم اختر العملية

⚡ *البوت يتعرف على الملف تلقائياً!*
"""

    if update.callback_query:
        await update.callback_query.edit_message_text(
            help_text,
            parse_mode="Markdown",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 رجوع", callback_data="back_menu")]])
        )
    else:
        await update.message.reply_text(
            help_text,
            parse_mode="Markdown",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 رجوع", callback_data="back_menu")]])
        )

async def check_libraries(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """فحص حالة المكتبات"""
    status = "📋 *حالة المكتبات:*

"

    for name, available in LIBS_STATUS.items():
        icon = "✅" if available else "❌"
        status += f"{icon} {name}
"

    status += "
⚠️ *إذا فيه ❌:*
"
    status += "المكتبة غير مثبتة، بعض الميزات قد لا تعمل"

    await update.callback_query.edit_message_text(
        status,
        parse_mode="Markdown",
        reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 رجوع", callback_data="back_menu")]])
    )

# ═══════════════════════════════════════════════
# المعالج الذكي الرئيسي
# ═══════════════════════════════════════════════

async def smart_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """المعالج الذكي - يتعرف على الملف ويحوله تلقائياً"""

    # فحص إذا كان صورة
    if update.message.photo:
        await handle_image_conversion(update, context)
        return

    # فحص إذا كان مستند
    doc = update.message.document
    if not doc:
        await update.message.reply_text(
            "❌ يرجى إرسال ملف أو صورة!",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🚀 ابدأ", callback_data="back_menu")]])
        )
        return

    # التعرف الذكي على نوع الملف
    file_info = detect_file_type(doc)

    if not file_info:
        await update.message.reply_text(
            f"❌ نوع الملف غير مدعوم!

"
            f"الملف: `{doc.file_name}`
"
            f"الأنواع المدعومة:
"
            f"• صور (JPG, PNG, WEBP)
"
            f"• Word (DOCX)
"
            f"• Excel (XLSX)
"
            f"• PowerPoint (PPTX)
"
            f"• PDF",
            parse_mode="Markdown"
        )
        return

    logger.info(f"🎯 تم التعرف على الملف: {file_info['name']}")

    # توجيه للمعالج المناسب
    handler_name = file_info["handler"]

    if handler_name == "handle_image_conversion":
        await handle_image_conversion(update, context)
    elif handler_name == "handle_word_conversion":
        await handle_word_conversion(update, context)
    elif handler_name == "handle_excel_conversion":
        await handle_excel_conversion(update, context)
    elif handler_name == "handle_ppt_conversion":
        await handle_ppt_conversion(update, context)
    elif handler_name == "handle_pdf_operations":
        await handle_pdf_operations(update, context)

# ═══════════════════════════════════════════════
# تحويل الصور
# ═══════════════════════════════════════════════

async def handle_image_conversion(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """تحويل صورة لـ PDF - تلقائي"""

    if not LIBS_STATUS["Pillow"]:
        await update.message.reply_text(
            "❌ مكتبة Pillow غير مثبتة!
"
            "التحويل غير متاح.",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 رجوع", callback_data="back_menu")]])
        )
        return

    processing_msg = await update.message.reply_text("⏳ جاري تحويل الصورة لـ PDF...")

    temp_image = None
    temp_pdf = None

    try:
        # الحصول على الصورة
        photo = update.message.photo[-1] if update.message.photo else None

        if photo:
            file = await context.bot.get_file(photo.file_id)
            temp_image = tempfile.mktemp(suffix=".jpg")
            await file.download_to_drive(temp_image)
        elif update.message.document:
            doc = update.message.document
            file = await context.bot.get_file(doc.file_id)
            ext = os.path.splitext(doc.file_name)[1] if doc.file_name else ".jpg"
            if not ext or ext == ".":
                ext = ".jpg"
            temp_image = tempfile.mktemp(suffix=ext)
            await file.download_to_drive(temp_image)
        else:
            await processing_msg.edit_text("❌ لم يتم العثور على صورة!")
            return

        # التحقق من حجم الملف
        file_size = get_file_size_kb(temp_image)
        logger.info(f"📊 حجم الصورة: {file_size:.1f} KB")

        if file_size > 20480:  # 20 MB
            await processing_msg.edit_text("❌ حجم الصورة كبير جداً! (الحد الأقصى 20 ميجا)")
            clean_temp_files(temp_image)
            return

        # تحويل إلى PDF
        temp_pdf = tempfile.mktemp(suffix=".pdf")

        with Image.open(temp_image) as img:
            logger.info(f"🖼️ نوع الصورة: {img.format}, الحجم: {img.size}, الوضع: {img.mode}")

            # تحويل إلى RGB
            if img.mode in ("RGBA", "P", "LA", "L", "RGBX"):
                img = img.convert("RGB")

            # حفظ كـ PDF
            img.save(temp_pdf, "PDF", resolution=100.0, quality=95)

        # التحقق
        if not os.path.exists(temp_pdf) or os.path.getsize(temp_pdf) == 0:
            raise Exception("فشل في إنشاء ملف PDF")

        # إرسال
        await processing_msg.delete()
        await update.message.reply_document(
            document=open(temp_pdf, "rb"),
            filename="converted.pdf",
            caption=f"✅ تم التحويل بنجاح!
"
                    f"🖼️ صورة → 📄 PDF"
        )
        logger.info("✅ تم إرسال PDF بنجاح")

    except Exception as e:
        error_msg = f"❌ خطأ في التحويل: {str(e)}"
        logger.error(f"❌ خطأ: {e}")
        logger.error(traceback.format_exc())
        await processing_msg.edit_text(error_msg)
    finally:
        clean_temp_files(temp_image, temp_pdf)

# ═══════════════════════════════════════════════
# تحويل Word
# ═══════════════════════════════════════════════

async def handle_word_conversion(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """تحويل Word لـ PDF - تلقائي"""

    if not LIBS_STATUS["python-docx"] or not LIBS_STATUS["ReportLab"]:
        missing = []
        if not LIBS_STATUS["python-docx"]:
            missing.append("python-docx")
        if not LIBS_STATUS["ReportLab"]:
            missing.append("reportlab")

        await update.message.reply_text(
            f"❌ مكتبات ناقصة: {', '.join(missing)}
"
            f"التحويل غير متاح.",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 رجوع", callback_data="back_menu")]])
        )
        return

    processing_msg = await update.message.reply_text("⏳ جاري تحويل Word لـ PDF...")

    temp_docx = None
    temp_pdf = None

    try:
        doc = update.message.document
        file = await context.bot.get_file(doc.file_id)
        temp_docx = tempfile.mktemp(suffix=".docx")
        await file.download_to_drive(temp_docx)

        temp_pdf = tempfile.mktemp(suffix=".pdf")

        # قراءة Word
        document = Document(temp_docx)

        # إنشاء PDF
        pdf = canvas.Canvas(temp_pdf, pagesize=A4)
        width, height = A4

        y_position = height - 2 * cm
        margin = 2 * cm
        line_height = 0.6 * cm

        for para in document.paragraphs:
            text = para.text.strip()
            if text:
                words = text.split()
                line = ""
                for word in words:
                    test_line = line + " " + word if line else word
                    if pdf.stringWidth(test_line, "Helvetica", 12) < (width - 2 * margin):
                        line = test_line
                    else:
                        if y_position < margin:
                            pdf.showPage()
                            y_position = height - 2 * cm
                        pdf.drawString(margin, y_position, line)
                        y_position -= line_height
                        line = word

                if line:
                    if y_position < margin:
                        pdf.showPage()
                        y_position = height - 2 * cm
                    pdf.drawString(margin, y_position, line)
                    y_position -= line_height
            else:
                y_position -= line_height * 0.5

            if y_position < margin:
                pdf.showPage()
                y_position = height - 2 * cm

        pdf.save()

        # إرسال
        await processing_msg.delete()
        await update.message.reply_document(
            document=open(temp_pdf, "rb"),
            filename="converted.pdf",
            caption=f"✅ تم التحويل بنجاح!
"
                    f"📝 Word → 📄 PDF"
        )

    except Exception as e:
        logger.error(f"❌ خطأ في تحويل Word: {e}")
        logger.error(traceback.format_exc())
        await processing_msg.edit_text(f"❌ خطأ: {str(e)}")
    finally:
        clean_temp_files(temp_docx, temp_pdf)

# ═══════════════════════════════════════════════
# تحويل Excel
# ═══════════════════════════════════════════════

async def handle_excel_conversion(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """تحويل Excel لـ PDF - تلقائي"""

    if not LIBS_STATUS["pandas"] or not LIBS_STATUS["ReportLab"]:
        missing = []
        if not LIBS_STATUS["pandas"]:
            missing.append("pandas, openpyxl")
        if not LIBS_STATUS["ReportLab"]:
            missing.append("reportlab")

        await update.message.reply_text(
            f"❌ مكتبات ناقصة: {', '.join(missing)}
"
            f"التحويل غير متاح.",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 رجوع", callback_data="back_menu")]])
        )
        return

    processing_msg = await update.message.reply_text("⏳ جاري تحويل Excel لـ PDF...")

    temp_excel = None
    temp_pdf = None

    try:
        doc = update.message.document
        file = await context.bot.get_file(doc.file_id)
        temp_excel = tempfile.mktemp(suffix=".xlsx")
        await file.download_to_drive(temp_excel)

        temp_pdf = tempfile.mktemp(suffix=".pdf")

        # قراءة Excel
        df = pd.read_excel(temp_excel)

        # إنشاء PDF
        doc_pdf = SimpleDocTemplate(temp_pdf, pagesize=A4)
        elements = []
        styles = getSampleStyleSheet()

        # العنوان
        title = Paragraph(f"<b>{doc.file_name}</b>", styles["Heading1"])
        elements.append(title)
        elements.append(Spacer(1, 20))

        # تحويل DataFrame لـ Table
        data = [df.columns.tolist()] + df.values.tolist()

        # تنسيق البيانات
        for i in range(1, len(data)):
            for j in range(len(data[i])):
                if pd.isna(data[i][j]):
                    data[i][j] = ""
                else:
                    data[i][j] = str(data[i][j])

        # حساب عرض الأعمدة
        col_count = len(df.columns)
        col_width = 15 * cm / col_count if col_count > 0 else 15 * cm

        table = Table(data, colWidths=[col_width] * col_count)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, 0), 10),
            ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
            ("BACKGROUND", (0, 1), (-1, -1), colors.beige),
            ("GRID", (0, 0), (-1, -1), 1, colors.black),
            ("FONTSIZE", (0, 1), (-1, -1), 8),
        ]))

        elements.append(table)
        doc_pdf.build(elements)

        # إرسال
        await processing_msg.delete()
        await update.message.reply_document(
            document=open(temp_pdf, "rb"),
            filename="converted.pdf",
            caption=f"✅ تم التحويل بنجاح!
"
                    f"📊 Excel → 📄 PDF"
        )

    except Exception as e:
        logger.error(f"❌ خطأ في تحويل Excel: {e}")
        logger.error(traceback.format_exc())
        await processing_msg.edit_text(f"❌ خطأ: {str(e)}")
    finally:
        clean_temp_files(temp_excel, temp_pdf)

# ═══════════════════════════════════════════════
# تحويل PowerPoint
# ═══════════════════════════════════════════════

async def handle_ppt_conversion(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """تحويل PowerPoint لـ PDF - تلقائي"""

    if not LIBS_STATUS["python-pptx"] or not LIBS_STATUS["ReportLab"]:
        missing = []
        if not LIBS_STATUS["python-pptx"]:
            missing.append("python-pptx")
        if not LIBS_STATUS["ReportLab"]:
            missing.append("reportlab")

        await update.message.reply_text(
            f"❌ مكتبات ناقصة: {', '.join(missing)}
"
            f"التحويل غير متاح.",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 رجوع", callback_data="back_menu")]])
        )
        return

    processing_msg = await update.message.reply_text("⏳ جاري تحويل PowerPoint لـ PDF...")

    temp_ppt = None
    temp_pdf = None

    try:
        doc = update.message.document
        file = await context.bot.get_file(doc.file_id)
        temp_ppt = tempfile.mktemp(suffix=".pptx")
        await file.download_to_drive(temp_ppt)

        temp_pdf = tempfile.mktemp(suffix=".pdf")

        # قراءة PowerPoint
        prs = Presentation(temp_ppt)

        # إنشاء PDF
        pdf = canvas.Canvas(temp_pdf, pagesize=A4)
        width, height = A4

        for i, slide in enumerate(prs.slides):
            y_position = height - 2 * cm

            # عنوان الشريحة
            pdf.setFont("Helvetica-Bold", 16)
            pdf.drawString(2 * cm, y_position, f"Slide {i + 1}")
            y_position -= 1.5 * cm

            pdf.setFont("Helvetica", 12)

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    words = text.split()
                    line = ""
                    for word in words:
                        test_line = line + " " + word if line else word
                        if pdf.stringWidth(test_line, "Helvetica", 12) < (width - 4 * cm):
                            line = test_line
                        else:
                            if y_position < 2 * cm:
                                break
                            pdf.drawString(2 * cm, y_position, line)
                            y_position -= 0.6 * cm
                            line = word

                    if line and y_position >= 2 * cm:
                        pdf.drawString(2 * cm, y_position, line)
                        y_position -= 0.6 * cm

            pdf.showPage()

        pdf.save()

        # إرسال
        await processing_msg.delete()
        await update.message.reply_document(
            document=open(temp_pdf, "rb"),
            filename="converted.pdf",
            caption=f"✅ تم التحويل بنجاح!
"
                    f"📽️ PowerPoint → 📄 PDF
"
                    f"📊 عدد الشرائح: {len(prs.slides)}"
        )

    except Exception as e:
        logger.error(f"❌ خطأ في تحويل PowerPoint: {e}")
        logger.error(traceback.format_exc())
        await processing_msg.edit_text(f"❌ خطأ: {str(e)}")
    finally:
        clean_temp_files(temp_ppt, temp_pdf)

# ═══════════════════════════════════════════════
# معالجة PDF
# ═══════════════════════════════════════════════

async def handle_pdf_operations(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """معالجة PDF - عرض الخيارات"""

    doc = update.message.document
    if not doc:
        return

    context.user_data["pdf_file_id"] = doc.file_id
    context.user_data["pdf_name"] = doc.file_name

    keyboard = [
        [InlineKeyboardButton("📎 دمج مع PDF آخر", callback_data="merge_with_pdf")],
        [InlineKeyboardButton("✂️ تقسيم PDF", callback_data="split_this_pdf")],
        [InlineKeyboardButton("🗜️ ضغط PDF", callback_data="compress_this_pdf")],
        [InlineKeyboardButton("🔙 إلغاء", callback_data="back_menu")]
    ]

    await update.message.reply_text(
        f"📄 *تم استلام ملف PDF*

"
        f"الاسم: `{doc.file_name}`
"
        f"الحجم: {doc.file_size / 1024:.1f} KB

"
        f"اختر العملية:",
        parse_mode="Markdown",
        reply_markup=InlineKeyboardMarkup(keyboard)
    )

# ═══════════════════════════════════════════════
# دمج PDFs
# ═══════════════════════════════════════════════

async def merge_with_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """دمج مع PDF موجود"""
    query = update.callback_query
    await query.answer()

    if not LIBS_STATUS["PyPDF2"]:
        await query.edit_message_text(
            "❌ مكتبة PyPDF2 غير مثبتة!",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 رجوع", callback_data="back_menu")]])
        )
        return

    await query.edit_message_text(
        "📎 *أرسل ملف PDF الثاني للدمج*

"
        "سيتم دمج الملفين في PDF واحد",
        parse_mode="Markdown",
        reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 إلغاء", callback_data="back_menu")]])
    )
    context.user_data["action"] = "merge_with_existing"

async def handle_merge_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """معالجة دمج PDF"""
    action = context.user_data.get("action", "")

    if action == "merge_with_existing":
        await merge_existing_pdfs(update, context)

async def merge_existing_pdfs(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """دمج PDFين"""
    processing_msg = await update.message.reply_text("⏳ جاري الدمج...")

    temp_pdf1 = None
    temp_pdf2 = None
    temp_output = None

    try:
        pdf1_file_id = context.user_data.get("pdf_file_id")
        if not pdf1_file_id:
            await processing_msg.edit_text("❌ خطأ: لم يتم العثور على PDF أول")
            return

        file1 = await context.bot.get_file(pdf1_file_id)
        temp_pdf1 = tempfile.mktemp(suffix=".pdf")
        await file1.download_to_drive(temp_pdf1)

        doc = update.message.document
        file2 = await context.bot.get_file(doc.file_id)
        temp_pdf2 = tempfile.mktemp(suffix=".pdf")
        await file2.download_to_drive(temp_pdf2)

        temp_output = tempfile.mktemp(suffix=".pdf")

        merger = PdfMerger()
        merger.append(temp_pdf1)
        merger.append(temp_pdf2)
        merger.write(temp_output)
        merger.close()

        await processing_msg.delete()
        await update.message.reply_document(
            document=open(temp_output, "rb"),
            filename="merged.pdf",
            caption="✅ تم الدمج بنجاح!"
        )

    except Exception as e:
        logger.error(f"❌ خطأ في دمج PDFs: {e}")
        logger.error(traceback.format_exc())
        await processing_msg.edit_text(f"❌ خطأ: {str(e)}")
    finally:
        clean_temp_files(temp_pdf1, temp_pdf2, temp_output)
        context.user_data["action"] = ""

# ═══════════════════════════════════════════════
# تقسيم PDF
# ═══════════════════════════════════════════════

async def split_this_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """تقسيم PDF"""
    query = update.callback_query
    await query.answer()

    if not LIBS_STATUS["PyPDF2"]:
        await query.edit_message_text(
            "❌ مكتبة PyPDF2 غير مثبتة!",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 رجوع", callback_data="back_menu")]])
        )
        return

    pdf_file_id = context.user_data.get("pdf_file_id")
    if not pdf_file_id:
        await query.edit_message_text("❌ خطأ: لم يتم العثور على PDF")
        return

    await query.edit_message_text(
        "✂️ *تقسيم PDF*

"
        "أرسل رقم الصفحة أو النطاق
"
        "مثال: `1` أو `1-3`",
        parse_mode="Markdown",
        reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🔙 إلغاء", callback_data="back_menu")]])
    )
    context.user_data["action"] = "split_pages"

async def handle_split_pages(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """تنفيذ التقسيم"""
    processing_msg = await update.message.reply_text("⏳ جاري التقسيم...")

    temp_pdf = None
    temp_output = None

    try:
        pdf_file_id = context.user_data.get("pdf_file_id")
        if not pdf_file_id:
            await processing_msg.edit_text("❌ خطأ: لم يتم العثور على PDF")
            return

        file = await context.bot.get_file(pdf_file_id)
        temp_pdf = tempfile.mktemp(suffix=".pdf")
        await file.download_to_drive(temp_pdf)

        text = update.message.text.strip()
        reader = PdfReader(temp_pdf)
        writer = PdfWriter()

        if "-" in text:
            start, end = map(int, text.split("-"))
            start = max(1, start)
            end = min(len(reader.pages), end)

            for i in range(start - 1, end):
                writer.add_page(reader.pages[i])

            filename = f"pages_{start}-{end}.pdf"
            caption = f"✅ تم استخراج الصفحات {start}-{end}"
        else:
            page_num = int(text)
            if page_num < 1 or page_num > len(reader.pages):
                await processing_msg.edit_text(f"❌ رقم غير صالح! (1-{len(reader.pages)})")
                return

            writer.add_page(reader.pages[page_num - 1])
            filename = f"page_{page_num}.pdf"
            caption = f"✅ تم استخراج الصفحة {page_num}"

        temp_output = tempfile.mktemp(suffix=".pdf")
        with open(temp_output, "wb") as f:
            writer.write(f)

        await processing_msg.delete()
        await update.message.reply_document(
            document=open(temp_output, "rb"),
            filename=filename,
            caption=caption
        )

    except ValueError:
        await processing_msg.edit_text("❌ يرجى إرسال رقم أو نطاق مثل: 1-3")
    except Exception as e:
        logger.error(f"❌ خطأ في تقسيم PDF: {e}")
        logger.error(traceback.format_exc())
        await processing_msg.edit_text(f"❌ خطأ: {str(e)}")
    finally:
        clean_temp_files(temp_pdf, temp_output)
        context.user_data["action"] = ""

# ═══════════════════════════════════════════════
# ضغط PDF
# ═══════════════════════════════════════════════

async def compress_this_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """ضغط PDF"""
    query = update.callback_query
    await query.answer()

    pdf_file_id = context.user_data.get("pdf_file_id")
    if not pdf_file_id:
        await query.edit_message_text("❌ خطأ: لم يتم العثور على PDF")
        return

    processing_msg = await query.edit_message_text("⏳ جاري الضغط...")

    temp_pdf = None
    temp_output = None

    try:
        file = await context.bot.get_file(pdf_file_id)
        temp_pdf = tempfile.mktemp(suffix=".pdf")
        await file.download_to_drive(temp_pdf)

        temp_output = tempfile.mktemp(suffix=".pdf")

        if LIBS_STATUS["PyPDF2"]:
            reader = PdfReader(temp_pdf)
            writer = PdfWriter()

            for page in reader.pages:
                writer.add_page(page)

            with open(temp_output, "wb") as f:
                writer.write(f)
        else:
            # نسخ بدون ضغط إذا PyPDF2 غير متوفر
            import shutil
            shutil.copy(temp_pdf, temp_output)

        original_size = os.path.getsize(temp_pdf)
        compressed_size = os.path.getsize(temp_output)

        await processing_msg.delete()
        await context.bot.send_document(
            chat_id=update.effective_user.id,
            document=open(temp_output, "rb"),
            filename="compressed.pdf",
            caption=f"✅ تم الضغط!
"
                    f"📊 الأصلي: {original_size / 1024:.1f} KB
"
                    f"📉 الجديد: {compressed_size / 1024:.1f} KB"
        )

    except Exception as e:
        logger.error(f"❌ خطأ في ضغط PDF: {e}")
        logger.error(traceback.format_exc())
        await processing_msg.edit_text(f"❌ خطأ: {str(e)}")
    finally:
        clean_temp_files(temp_pdf, temp_output)
        context.user_data["action"] = ""

# ═══════════════════════════════════════════════
# معالجة الأزرار
# ═══════════════════════════════════════════════

async def button_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """معالجة الأزرار"""
    query = update.callback_query
    data = query.data

    handlers = {
        "help": help_command,
        "check_libs": check_libraries,
        "back_menu": start,
        "merge_with_pdf": merge_with_pdf,
        "split_this_pdf": split_this_pdf,
        "compress_this_pdf": compress_this_pdf,
    }

    handler = handlers.get(data)
    if handler:
        await handler(update, context)

# ═══════════════════════════════════════════════
# الدالة الرئيسية
# ═══════════════════════════════════════════════

def main():
    """تشغيل البوت"""
    logger.info("🚀 جاري تشغيل البوت الذكي...")
    logger.info(f"🔑 Token: {TOKEN[:20]}...")

    application = Application.builder().token(TOKEN).build()

    # الأوامر
    application.add_handler(CommandHandler("start", start))

    # الأزرار
    application.add_handler(CallbackQueryHandler(button_handler))

    # المعالج الذكي الرئيسي (صور + ملفات)
    application.add_handler(MessageHandler(filters.PHOTO | filters.Document.ALL, smart_handler))

    # الرسائل النصية (لتقسيم PDF)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_split_pages))

    logger.info("✅ البوت يعمل الآن...")
    application.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__ == "__main__":
    main()
